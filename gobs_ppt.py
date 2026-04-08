"""
GOBS Platform PPT Generator
Generates a comprehensive pitch deck covering:
  1. Where We Are
  2. Where We Aim to Go
  3. How to Achieve (Roadmap)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x0F, 0x1A)   # deep navy
BG_CARD       = RGBColor(0x16, 0x1B, 0x2E)   # card bg
ACCENT        = RGBColor(0x5B, 0x8D, 0xFF)   # primary blue
ACCENT2       = RGBColor(0x7C, 0x3A, 0xFF)   # purple
ACCENT_GREEN  = RGBColor(0x2E, 0xCC, 0x8B)   # green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)   # orange
ACCENT_RED    = RGBColor(0xFF, 0x4D, 0x6D)   # red
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT    = RGBColor(0xB0, 0xB8, 0xD0)
GRAY_MED      = RGBColor(0x60, 0x6B, 0x8A)
YELLOW        = RGBColor(0xFF, 0xD6, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, alpha=None, radius=False):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def tag(slide, text, l, t, color=ACCENT):
    """Small pill tag."""
    rect(slide, l, t, len(text) * 0.095 + 0.2, 0.28, color)
    txt(slide, text, l + 0.08, t + 0.02, len(text) * 0.11 + 0.1, 0.28,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def divider(slide, t, color=ACCENT, l=0.5, w=12.33):
    rect(slide, l, t, w, 0.03, color)


def section_label(slide, text, t=0.22):
    txt(slide, text, 0.5, t, 4, 0.35, size=10, bold=True,
        color=ACCENT, align=PP_ALIGN.LEFT)


def slide_title(slide, title, subtitle=None, t=0.55):
    txt(slide, title, 0.5, t, 12.33, 0.6,
        size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.5, t + 0.62, 12.33, 0.4,
            size=16, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)


def card(slide, l, t, w, h, color=BG_CARD):
    return rect(slide, l, t, w, h, color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# Gradient-feel top bar
rect(s, 0, 0, 13.33, 0.06, ACCENT)

# Large title
txt(s, "GOBS", 0.6, 1.8, 10, 1.6,
    size=96, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Gradient accent line under logo
rect(s, 0.6, 3.35, 5, 0.07, ACCENT)

# Subtitle
txt(s, "Game Operations & Broadcasting System",
    0.62, 3.55, 10, 0.5,
    size=20, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

txt(s, "Product Vision · Architecture · Roadmap",
    0.62, 4.1, 10, 0.4,
    size=14, color=GRAY_MED, align=PP_ALIGN.LEFT)

# Version badge
rect(s, 0.6, 4.7, 1.5, 0.32, ACCENT2)
txt(s, "v2.0  |  Apr 2026", 0.65, 4.73, 1.4, 0.28,
    size=10, bold=True, color=WHITE)

# Right side visual - abstract circles
rect(s, 8.5, 1.0, 4.0, 5.5, RGBColor(0x16, 0x1B, 0x2E))
txt(s, "🚀", 9.8, 2.5, 2, 1.5, size=80, align=PP_ALIGN.CENTER)

# Bottom bar
rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))
txt(s, "Confidential  ·  Internal Use Only", 0.5, 7.22, 12, 0.25,
    size=9, color=GRAY_MED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT)

section_label(s, "AGENDA")
slide_title(s, "Today's Agenda", t=0.45)
divider(s, 1.15)

items = [
    ("01", "Where We Are",       "Current capabilities of GOBS H5 platform",              ACCENT),
    ("02", "Where We Aim to Go", "Vision: AI-driven game publishing flywheel",              ACCENT2),
    ("03", "How to Achieve",     "Architecture, core modules & delivery roadmap",           ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(items):
    t0 = 1.5 + i * 1.7
    card(s, 0.5, t0, 12.33, 1.5)
    rect(s, 0.5, t0, 0.08, 1.5, color)
    txt(s, num,   0.75, t0 + 0.15, 0.8, 0.5, size=28, bold=True, color=color)
    txt(s, title, 1.65, t0 + 0.18, 8, 0.45, size=22, bold=True, color=WHITE)
    txt(s, desc,  1.65, t0 + 0.72, 9, 0.35, size=13, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Section Divider: WHERE WE ARE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 7.5, BG_DARK)
rect(s, 0, 0, 0.12, 7.5, ACCENT)

txt(s, "01", 1.0, 1.5, 4, 2.0, size=120, bold=True,
    color=RGBColor(0x1A, 0x20, 0x35))

txt(s, "WHERE\nWE ARE", 1.0, 2.2, 10, 3.0,
    size=56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "A look at what GOBS H5 can do today —\nand where the gaps are.",
    1.0, 5.4, 9, 0.8,
    size=16, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

rect(s, 1.0, 5.2, 3.5, 0.05, ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Current Platform Overview
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT)

section_label(s, "01  WHERE WE ARE")
slide_title(s, "GOBS H5 — Current Platform", "What exists and works today", t=0.45)
divider(s, 1.15)

modules = [
    ("🎬", "Content Production",
     ["Studio: text/reference-based video gen", "Seedance 2.0 integration (ViralDance)", "Multi-shot storyboard with style lock", "Character portrait editor + look tree"],
     ACCENT),
    ("🧑‍🎨", "Character System",
     ["Character card with look tree (variants)", "Wardrobe panel: base image + state costumes", "Character library: save / reuse across projects", "Scene & prop card image generation"],
     ACCENT2),
    ("📤", "Distribution",
     ["GeeLark account matrix integration", "Batch overnight publishing", "TikTok / Instagram multi-account support", "Task board with status tracking"],
     ACCENT_GREEN),
    ("🗂️", "Production Wizard",
     ["Advanced production flow (高级制片)", "Template market & save-as-template", "Drive material picker", "Batch job queue management"],
     ACCENT_ORANGE),
]

col_w = 3.0
for i, (icon, title, bullets, color) in enumerate(modules):
    col = i % 2
    row = i // 2
    l = 0.5 + col * 6.4
    t = 1.4 + row * 2.8
    card(s, l, t, 6.1, 2.55)
    rect(s, l, t, 6.1, 0.06, color)
    txt(s, icon + "  " + title, l + 0.2, t + 0.15, 5.5, 0.45,
        size=14, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, "·  " + b, l + 0.25, t + 0.6 + j * 0.45, 5.5, 0.4,
            size=11, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Current Gaps (Pain Points)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT)

section_label(s, "01  WHERE WE ARE  ·  GAPS")
slide_title(s, "What's Missing Today", "The flywheel is broken — no closed loop", t=0.45)
divider(s, 1.15)

# Flow diagram: Production → Distribution → ??? → (no feedback)
flow_items = ["Content\nProduction", "Distribution\n(GeeLark)", "Performance\nData", "Insights &\nDecisions"]
flow_colors = [ACCENT, ACCENT2, ACCENT_ORANGE, ACCENT_GREEN]
flow_x = [0.6, 3.4, 6.2, 9.0]

for i, (label, color, x) in enumerate(zip(flow_items, flow_colors, flow_x)):
    is_missing = i >= 2
    bg_c = RGBColor(0x2A, 0x10, 0x10) if is_missing else BG_CARD
    card(s, x, 1.4, 2.5, 1.3, color=bg_c)
    if is_missing:
        rect(s, x, 1.4, 2.5, 0.05, ACCENT_RED)
    else:
        rect(s, x, 1.4, 2.5, 0.05, color)
    align = PP_ALIGN.CENTER
    c = GRAY_MED if is_missing else WHITE
    txt(s, label, x + 0.1, t+0.35, 2.3, 0.7, size=13, bold=not is_missing,
        color=c, align=align)
    if is_missing:
        txt(s, "❌ Missing", x + 0.5, 1.45, 1.5, 0.3,
            size=9, bold=True, color=ACCENT_RED)
    if i < 3:
        txt(s, "→", x + 2.52, 1.8, 0.3, 0.5, size=20, bold=True,
            color=GRAY_MED, align=PP_ALIGN.CENTER)

# Gap descriptions
gaps = [
    ("No Data Feedback Loop",
     "After publishing, performance data (views, completion rate, engagement) never returns to the platform. There is no causal link between 'what we made' and 'how it performed'.",
     ACCENT_RED),
    ("No Intelligence Layer",
     "Zero AI-driven insight. No attribution analysis. No understanding of which content type, character, or timing works best. All decisions are manual and gut-feel.",
     ACCENT_ORANGE),
    ("No Actionable Recommendations",
     "Even if data existed, there is no mechanism to translate signals into concrete actions: produce this video, engage this KOL, fix this game system.",
     YELLOW),
    ("No Self-Improvement",
     "The platform cannot learn from past actions. Every campaign starts from scratch. No compounding intelligence.",
     ACCENT2),
]

for i, (title, desc, color) in enumerate(gaps):
    l = 0.5 + (i % 2) * 6.4
    t = 3.0 + (i // 2) * 2.0
    card(s, l, t, 6.1, 1.75)
    rect(s, l, t, 0.06, 1.75, color)
    txt(s, title, l + 0.2, t + 0.15, 5.7, 0.35, size=13, bold=True, color=color)
    txt(s, desc,  l + 0.2, t + 0.5,  5.7, 1.1,  size=10, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section Divider: WHERE WE AIM TO GO
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT2)

txt(s, "02", 1.0, 1.5, 4, 2.0, size=120, bold=True,
    color=RGBColor(0x1A, 0x20, 0x35))

txt(s, "WHERE\nWE AIM\nTO GO", 1.0, 1.8, 10, 3.5,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "An AI-powered publishing flywheel that learns,\nadapts, and eventually runs itself.",
    1.0, 5.5, 9, 0.8,
    size=16, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

rect(s, 1.0, 5.35, 5, 0.05, ACCENT2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Vision
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT2)

section_label(s, "02  WHERE WE AIM TO GO")
slide_title(s, "The GOBS Vision", "Bind your game. Let AI run your publishing.", t=0.45)
divider(s, 1.15, color=ACCENT2)

txt(s,
    "A game operator binds their game once. From that moment, GOBS continuously senses internal and "
    "external signals, generates AI-driven insights, recommends and executes actions — "
    "then learns from every outcome. Over time, the platform runs the publishing flywheel autonomously.",
    0.5, 1.3, 12.33, 0.8, size=13, color=GRAY_LIGHT)

# Flywheel circles (text-based)
flywheel = [
    ("📡", "Sense",    "Collect game KPIs,\ncontent performance,\nbuy data, sentiment, KOL",   0.6,  2.3),
    ("🧠", "Insight",  "AI attribution,\nstage-aware advice,\nROI-ranked actions",              3.6,  2.3),
    ("⚡", "Act",      "Agent executes:\nvideo, distribution,\nKOL, crisis, PRD",              6.6,  2.3),
    ("📊", "Track",    "ROI attribution,\ncausal link:\naction → outcome",                     9.6,  2.3),
]

for icon, title, desc, l, t in flywheel:
    card(s, l, t, 2.75, 2.9)
    txt(s, icon, l + 1.0, t + 0.2, 0.8, 0.6, size=28, align=PP_ALIGN.CENTER)
    txt(s, title, l + 0.1, t + 0.85, 2.5, 0.45, size=16, bold=True,
        color=ACCENT2, align=PP_ALIGN.CENTER)
    txt(s, desc,  l + 0.1, t + 1.35, 2.5, 1.3,  size=11, color=GRAY_LIGHT,
        align=PP_ALIGN.CENTER)

# Arrows between boxes
for ax in [3.37, 6.37, 9.37]:
    txt(s, "→", ax, 3.4, 0.28, 0.5, size=18, bold=True,
        color=ACCENT2, align=PP_ALIGN.CENTER)

# Feedback arrow label
txt(s, "↩  Self-learning feedback loop", 3.8, 5.35, 5.5, 0.35,
    size=11, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# Three modes
modes = [
    ("Assisted Mode",    "AI recommends,\nhuman approves all",  "Month 0–3",  ACCENT),
    ("Semi-Auto Mode",   "Low-risk actions auto-execute,\nhigh-risk human review", "Month 3–6",  ACCENT2),
    ("Full-Auto Mode",   "Set direction & guardrails,\nplatform self-operates",   "Month 6+",   ACCENT_GREEN),
]
for i, (m, d, when, color) in enumerate(modes):
    l = 0.5 + i * 4.27
    card(s, l, 5.75, 4.0, 1.55)
    rect(s, l, 5.75, 4.0, 0.05, color)
    txt(s, m,    l + 0.2, 5.85, 3.5, 0.38, size=13, bold=True, color=color)
    txt(s, d,    l + 0.2, 6.25, 3.5, 0.65, size=10, color=GRAY_LIGHT)
    txt(s, when, l + 0.2, 6.9,  3.5, 0.28, size=9,  color=GRAY_MED)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Target Architecture (Full Stack)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT2)

section_label(s, "02  WHERE WE AIM TO GO  ·  ARCHITECTURE")
slide_title(s, "Full Platform Architecture", "8-layer stack from data to autonomous execution", t=0.45)
divider(s, 1.15, color=ACCENT2)

layers = [
    ("L0", "Onboarding",           "Game bind · Publishing stage config · Compliance setup · Knowledge base init",           ACCENT),
    ("L1", "Signal Collector",     "In-game KPIs · Content performance · Buy data (Meta/TikTok Ads) · Sentiment · KOL · Competitor", ACCENT_GREEN),
    ("L2", "Insight Engine",       "Attribution analysis · Stage-aware recommendations · ROI-ranked action list",            ACCENT2),
    ("L3", "Compliance Gate",      "Content rules · Regional filters · Risk scoring (auto / manual routing)",                ACCENT_ORANGE),
    ("L4", "Action Agents",        "Content Agent · Distribution Agent · KOL Agent · Crisis Agent · Product Agent · Ads Agent", ACCENT),
    ("L5", "Human Review Gate",    "Official publish · KOL deals · Crisis PR · Budget changes  →  always human",            ACCENT_RED),
    ("L6", "ROI Tracker",          "Cost per action · Incremental value attribution · Weekly ROI leaderboard",               YELLOW),
    ("L7", "Memory & Learning",    "Action archive · Strategy weight update · Cross-game knowledge transfer",                ACCENT2),
]

row_h = 0.68
for i, (code, name, desc, color) in enumerate(layers):
    t = 1.35 + i * row_h
    card(s, 0.5, t, 12.33, row_h - 0.05)
    rect(s, 0.5, t, 0.07, row_h - 0.05, color)
    txt(s, code, 0.65, t + 0.1, 0.55, 0.45, size=9, bold=True, color=color)
    txt(s, name, 1.25, t + 0.06, 2.3, 0.3, size=12, bold=True, color=WHITE)
    txt(s, desc, 1.25, t + 0.36, 11.2, 0.28, size=10, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section Divider: HOW TO ACHIEVE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_GREEN)

txt(s, "03", 1.0, 1.5, 4, 2.0, size=120, bold=True,
    color=RGBColor(0x1A, 0x20, 0x35))

txt(s, "HOW TO\nACHIEVE", 1.0, 2.0, 10, 2.5,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "Core modules, feature breakdown,\nand a time-bound delivery roadmap.",
    1.0, 4.7, 9, 0.8,
    size=16, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

rect(s, 1.0, 4.55, 4.5, 0.05, ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Phase 1: Foundation (Month 1)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT_GREEN)

section_label(s, "03  HOW TO ACHIEVE  ·  PHASE 1")
slide_title(s, "Phase 1 — Foundation", "Month 1  |  Goal: Close the data loop", t=0.45)
divider(s, 1.15, color=ACCENT_GREEN)

tag(s, "Month 1", 0.5, 1.22, color=ACCENT_GREEN)

modules_p1 = [
    ("Game Onboarding Module", [
        ("Game bind: name, region, platform, team", "Week 1"),
        ("Publishing stage configuration (soft launch / burst / long-tail)", "Week 1"),
        ("Official social account OAuth (TikTok, YouTube, IG)", "Week 2"),
        ("Compliance config: regional ratings, content red lines", "Week 2"),
        ("Knowledge base initialization (Wiki, character lore, past content)", "Week 3"),
    ], ACCENT),
    ("Data Feedback Loop (L1 MVP)", [
        ("Content performance data ingestion (official account API)", "Week 2"),
        ("Video → outcome causal record (template / character / time / platform)", "Week 2"),
        ("Performance dashboard: daily metrics overview", "Week 3"),
        ("In-game data stub (DAU/MAU manual input or API)", "Week 3–4"),
        ("Buy-side data connection (TikTok Ads / Meta Ads API)", "Week 4"),
    ], ACCENT_GREEN),
    ("Content Production Enhancements", [
        ("Character library: save-to-library from character card  ✅ Done", "Done"),
        ("Wardrobe base image: auto-inherit active look  ✅ Done", "Done"),
        ("Batch variant production (3 variants per brief, A/B)", "Week 3"),
        ("Quality auto-scoring (face integrity / motion smoothness)", "Week 4"),
        ("Template auto-save from top-performing videos", "Week 4"),
    ], ACCENT2),
]

col_positions = [0.5, 4.7, 8.9]
for i, (title, items, color) in enumerate(modules_p1):
    l = col_positions[i]
    card(s, l, 1.65, 4.0, 5.4)
    rect(s, l, 1.65, 4.0, 0.06, color)
    txt(s, title, l + 0.15, 1.78, 3.7, 0.4, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        t0 = 2.25 + j * 0.95
        done = "Done" in when
        c = ACCENT_GREEN if done else GRAY_LIGHT
        txt(s, "✓ " if done else "·  ", l + 0.1, t0, 0.3, 0.35, size=10, color=c)
        txt(s, feat, l + 0.3, t0, 3.2, 0.45, size=10, color=c)
        txt(s, when, l + 0.3, t0 + 0.45, 3.2, 0.28, size=8,
            color=ACCENT_GREEN if done else GRAY_MED)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Phase 2: Closed Loop (Month 2–3)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT2)

section_label(s, "03  HOW TO ACHIEVE  ·  PHASE 2")
slide_title(s, "Phase 2 — Closed Loop", "Month 2–3  |  Goal: Insights drive actions", t=0.45)
divider(s, 1.15, color=ACCENT2)

tag(s, "Month 2–3", 0.5, 1.22, color=ACCENT2)

modules_p2 = [
    ("Insight Engine (L2)", [
        ("Sentiment analysis pipeline (Reddit / TikTok comments)", "Month 2"),
        ("Content attribution: character / template / timing analysis", "Month 2"),
        ("Publishing stage context: same signal → different advice per stage", "Month 2"),
        ("Competitor monitoring: content trends, buy-side signals, version timing", "Month 2–3"),
        ("ROI baseline: cost per action type vs incremental value", "Month 3"),
        ("Daily insight report + priority-ranked action list", "Month 3"),
    ], ACCENT2),
    ("Action Recommendation Layer", [
        ("Action brief generation (structured: trigger / action / expected outcome)", "Month 2"),
        ("Human approval workflow (confirm → queue → agent executes)", "Month 2"),
        ("KOL database: tiered profiles (tier 1/2/3), audience, pricing, history", "Month 2–3"),
        ("KOL opportunity detection (content gap, collab invite draft)", "Month 3"),
        ("Crisis detection: negative sentiment threshold → crisis mode trigger", "Month 3"),
        ("Crisis mode: auto-pause all publishing, draft official response", "Month 3"),
    ], ACCENT_ORANGE),
    ("Distribution & Compliance", [
        ("Regional content variant (subtitle / dubbing auto-gen per region)", "Month 2"),
        ("Compliance gate: content rule check before every publish action", "Month 2"),
        ("Risk scoring: auto-approve / flag / block by risk level", "Month 2"),
        ("Optimal publish time recommendation (ML from historical data)", "Month 3"),
        ("Multi-platform simultaneous publish with platform-specific formatting", "Month 3"),
    ], ACCENT_GREEN),
]

col_positions = [0.5, 4.7, 8.9]
for i, (title, items, color) in enumerate(modules_p2):
    l = col_positions[i]
    card(s, l, 1.65, 4.0, 5.4)
    rect(s, l, 1.65, 4.0, 0.06, color)
    txt(s, title, l + 0.15, 1.78, 3.7, 0.4, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        t0 = 2.28 + j * 0.83
        txt(s, "·  " + feat, l + 0.15, t0, 3.65, 0.45, size=10, color=GRAY_LIGHT)
        txt(s, when, l + 0.25, t0 + 0.45, 3.5, 0.25, size=8, color=GRAY_MED)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Phase 3: Intelligence (Month 4–6)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT_ORANGE)

section_label(s, "03  HOW TO ACHIEVE  ·  PHASE 3")
slide_title(s, "Phase 3 — Intelligence", "Month 4–6  |  Goal: Self-learning, semi-autonomous", t=0.45)
divider(s, 1.15, color=ACCENT_ORANGE)

tag(s, "Month 4–6", 0.5, 1.22, color=ACCENT_ORANGE)

modules_p3 = [
    ("Action Agents (Full Suite)", [
        ("Content Agent: prompt gen → Seedance → quality score → review queue", "Month 4"),
        ("Distribution Agent: optimal-time auto-publish, matrix + official", "Month 4"),
        ("KOL Agent: collab tracking, invite draft, performance attribution", "Month 4"),
        ("Ads Agent: creative brief from top performers, budget realloc suggestion", "Month 4–5"),
        ("Product Agent: player feedback → root cause → PRD draft → Jira link", "Month 5"),
        ("Crisis Agent: sentiment watch → crisis mode → response drafting", "Month 5"),
    ], ACCENT_ORANGE),
    ("Memory & Self-Learning (L7)", [
        ("Action archive: trigger → execution → ROI outcome full causal record", "Month 4"),
        ("Weekly strategy weight update: high-ROI action types promoted", "Month 5"),
        ("Agent default parameter tuning (timing, style, character preference)", "Month 5"),
        ("Knowledge base auto-update on new game version releases", "Month 5"),
        ("Negative feedback learning: rejected actions reduce recommendation weight", "Month 5–6"),
        ("Industry benchmark DB: normal CPM / completion rate by genre", "Month 6"),
    ], ACCENT2),
    ("ROI Tracker & Reporting (L6)", [
        ("Per-action cost model (tool cost + estimated human-hour equivalent)", "Month 4"),
        ("Incremental value attribution (DAU / revenue lift from action)", "Month 4–5"),
        ("Weekly ROI leaderboard: which action types deliver most value", "Month 5"),
        ("Budget allocation recommendation: shift spend to highest-ROI actions", "Month 5"),
        ("Semi-auto mode: low-risk actions auto-execute without approval", "Month 6"),
        ("Cross-game learning: effective tactics from game A inform game B", "Month 6"),
    ], ACCENT_GREEN),
]

col_positions = [0.5, 4.7, 8.9]
for i, (title, items, color) in enumerate(modules_p3):
    l = col_positions[i]
    card(s, l, 1.65, 4.0, 5.4)
    rect(s, l, 1.65, 4.0, 0.06, color)
    txt(s, title, l + 0.15, 1.78, 3.7, 0.4, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        t0 = 2.28 + j * 0.83
        txt(s, "·  " + feat, l + 0.15, t0, 3.65, 0.45, size=10, color=GRAY_LIGHT)
        txt(s, when, l + 0.25, t0 + 0.45, 3.5, 0.25, size=8, color=GRAY_MED)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Phase 4: Flywheel (Month 6+)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT_GREEN)

section_label(s, "03  HOW TO ACHIEVE  ·  PHASE 4")
slide_title(s, "Phase 4 — Flywheel", "Month 6+  |  Goal: Platform runs itself", t=0.45)
divider(s, 1.15, color=ACCENT_GREEN)

tag(s, "Month 6+", 0.5, 1.22, color=ACCENT_GREEN)

items_p4 = [
    ("Full-Auto Mode",
     "Low-risk actions execute automatically. Humans set direction + guardrails. Weekly review replaces daily approval.",
     ACCENT_GREEN, 0.5, 1.65),
    ("3D Optimization Matrix",
     "Account persona × content type × publish timing — continuously optimized by the learning layer for each platform and region.",
     ACCENT2, 0.5, 3.25),
    ("Cross-Game Intelligence",
     "Successful tactics discovered for Game A surface as recommendations for Game B. Platform builds genre-level playbooks over time.",
     ACCENT_ORANGE, 0.5, 4.85),
    ("Autonomous Publishing Flywheel",
     "Sense → Insight → Act → Track → Learn → repeat. Human role shifts from 'operator' to 'director'. 10-person team operates at 100-person scale.",
     ACCENT, 7.0, 1.65),
    ("Platform-Level Intelligence",
     "Aggregated anonymized data across all games builds industry benchmarks. GOBS becomes the reference for what 'good' publishing looks like in each genre.",
     ACCENT2, 7.0, 3.25),
    ("What Stays Human (Always)",
     "Official brand publishing approval · KOL contract sign-off · Crisis PR final response · Major budget decisions · Product roadmap direction",
     ACCENT_RED, 7.0, 4.85),
]

for title, desc, color, l, t in items_p4:
    card(s, l, t, 5.8, 1.45)
    rect(s, l, t, 0.07, 1.45, color)
    txt(s, title, l + 0.2, t + 0.12, 5.3, 0.35, size=13, bold=True, color=color)
    txt(s, desc,  l + 0.2, t + 0.5,  5.3, 0.85, size=10, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Master Roadmap Timeline
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT)

section_label(s, "03  HOW TO ACHIEVE  ·  ROADMAP")
slide_title(s, "Master Roadmap", "Feature delivery timeline by phase", t=0.45)
divider(s, 1.15)

# Timeline header
phases = [
    ("Phase 1\nMonth 1",    ACCENT,        0.5),
    ("Phase 2\nMonth 2–3",  ACCENT2,       3.85),
    ("Phase 3\nMonth 4–6",  ACCENT_ORANGE, 7.2),
    ("Phase 4\nMonth 6+",   ACCENT_GREEN,  10.55),
]
for label, color, l in phases:
    rect(s, l, 1.35, 3.0, 0.55, color)
    txt(s, label, l + 0.1, 1.38, 2.8, 0.52, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

# Rows
rows = [
    ("Data Layer",      ["Game bind + KB init", "Sentiment + competitor",           "Buy-side deep analytics",        "Cross-game benchmark DB"]),
    ("Insight Engine",  ["Perf dashboard",      "Attribution + stage-aware advice",  "ROI-ranked actions",             "3D optimization matrix"]),
    ("Compliance",      ["Config setup",        "Compliance gate + risk scoring",     "—",                              "—"]),
    ("Action Agents",   ["Content Agent MVP",   "Distribution + KOL Agent",          "Full agent suite + crisis",      "Full-auto mode"]),
    ("Human Gate",      ["All actions",         "High-risk only",                    "High-risk only",                 "Brand / budget / crisis"]),
    ("ROI Tracker",     ["—",                   "Cost model baseline",               "ROI leaderboard + realloc",      "—"]),
    ("Learning Layer",  ["—",                   "—",                                 "Action archive + weight update", "Cross-game transfer"]),
]

row_colors = [ACCENT, ACCENT2, ACCENT_ORANGE, ACCENT_GREEN]

for i, (module, cells) in enumerate(rows):
    t = 2.0 + i * 0.72
    # Module label
    rect(s, 0.0, t, 0.5, 0.68, RGBColor(0x10, 0x14, 0x22))
    # Row bg alternating
    bg_c = BG_CARD if i % 2 == 0 else RGBColor(0x12, 0x17, 0x28)
    for j, cell in enumerate(cells):
        l = 0.5 + j * 3.35
        rect(s, l, t, 3.32, 0.68, bg_c)
        if cell != "—":
            txt(s, "·  " + cell, l + 0.1, t + 0.15, 3.1, 0.5, size=9.5,
                color=GRAY_LIGHT if j < 2 else WHITE)
        else:
            txt(s, "—", l + 1.3, t + 0.2, 0.5, 0.3, size=9, color=GRAY_MED)

    # Left module label (rotated text workaround — just left-align small box)
    txt(s, module, 0.02, t + 0.2, 0.48, 0.3, size=7.5, bold=True,
        color=ACCENT, align=PP_ALIGN.LEFT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Publishing Guardrails (What Stays Human)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT_RED)

section_label(s, "03  HOW TO ACHIEVE  ·  GUARDRAILS")
slide_title(s, "Automation vs Human Judgment", "Not everything should be automated — knowing the line is critical", t=0.45)
divider(s, 1.15, color=ACCENT_RED)

# Auto column
card(s, 0.5, 1.35, 5.8, 5.7)
rect(s, 0.5, 1.35, 5.8, 0.06, ACCENT_GREEN)
txt(s, "✅  Can Be Fully Automated", 0.65, 1.48, 5.4, 0.38,
    size=14, bold=True, color=ACCENT_GREEN)

auto_items = [
    "Data collection and daily dashboard update",
    "Attribution analysis and insight report generation",
    "Video production (enters review queue — not direct publish)",
    "Matrix / micro-account distribution",
    "Internal reports and PRD draft generation",
    "Quality auto-scoring and auto-regeneration",
    "Competitor and sentiment monitoring",
    "Optimal publish time recommendation",
]
for i, item in enumerate(auto_items):
    txt(s, "·  " + item, 0.65, 1.95 + i * 0.58, 5.5, 0.5, size=11, color=GRAY_LIGHT)

# Human column
card(s, 6.8, 1.35, 6.0, 5.7)
rect(s, 6.8, 1.35, 6.0, 0.06, ACCENT_RED)
txt(s, "🚫  Always Requires Human Approval", 6.95, 1.48, 5.6, 0.38,
    size=14, bold=True, color=ACCENT_RED)

human_items = [
    ("Official brand channel publishing",     "Brand risk — errors are public and permanent"),
    ("KOL partnership confirmation",           "Commercial relationship — needs human judgment"),
    ("Crisis PR response",                     "Reputational stakes — wrong tone = bigger crisis"),
    ("Cross-region sensitive content",         "Cultural nuance — AI cannot fully assess"),
    ("Major budget reallocation (>$X)",        "Financial risk — needs accountability"),
    ("Product roadmap / PRD submission",       "Business impact — requires stakeholder alignment"),
]
for i, (action, reason) in enumerate(human_items):
    t0 = 1.95 + i * 0.83
    txt(s, "⚠  " + action, 6.95, t0,        5.6, 0.35, size=11, bold=True, color=WHITE)
    txt(s, reason,          6.95, t0 + 0.37, 5.6, 0.3,  size=9,  color=GRAY_MED)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Closing / Call to Action
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.06, ACCENT)

txt(s, "The publishing flywheel\nstarts with one decision:", 1.0, 1.2, 10, 1.6,
    size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "Close the data loop.", 1.0, 2.9, 10, 0.9,
    size=48, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

rect(s, 1.0, 3.85, 6, 0.06, ACCENT)

txt(s,
    "Everything else — insight, agents, automation, self-learning — follows naturally "
    "once we can connect 'what we made' to 'what happened'. "
    "Phase 1 is the foundation. Let's start there.",
    1.0, 4.0, 11, 0.8, size=14, color=GRAY_LIGHT)

# Key priorities
priorities = [
    ("1st", "Data feedback loop (content → performance)", ACCENT),
    ("2nd", "Insight engine MVP with stage context",       ACCENT2),
    ("3rd", "Action agents + compliance gate",             ACCENT_GREEN),
    ("4th", "Self-learning + ROI tracker",                 ACCENT_ORANGE),
]
for i, (rank, desc, color) in enumerate(priorities):
    l = 1.0 + i * 3.0
    card(s, l, 5.0, 2.75, 1.6)
    rect(s, l, 5.0, 2.75, 0.06, color)
    txt(s, rank, l + 0.15, 5.15, 0.6, 0.45, size=18, bold=True, color=color)
    txt(s, desc, l + 0.15, 5.65, 2.45, 0.8, size=10, color=GRAY_LIGHT)

rect(s, 0, 7.2, 13.33, 0.3, RGBColor(0x10, 0x14, 0x22))
txt(s, "GOBS v2.0  ·  Apr 2026  ·  Confidential", 0.5, 7.22, 12, 0.25,
    size=9, color=GRAY_MED, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/openclaw/.openclaw/workspace/gobs_review/GOBS_Platform_Deck_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
