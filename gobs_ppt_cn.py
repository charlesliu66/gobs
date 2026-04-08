"""
GOBS Platform PPT — 中文版
解决布局挤压问题：所有元素严格计算坐标，不重叠
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 颜色 ──────────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0F, 0x1A)
CARD      = RGBColor(0x16, 0x1B, 0x2E)
CARD2     = RGBColor(0x12, 0x17, 0x26)
BLUE      = RGBColor(0x5B, 0x8D, 0xFF)
PURPLE    = RGBColor(0x7C, 0x3A, 0xFF)
GREEN     = RGBColor(0x2E, 0xCC, 0x8B)
ORANGE    = RGBColor(0xFF, 0x8C, 0x42)
RED       = RGBColor(0xFF, 0x4D, 0x6D)
YELLOW    = RGBColor(0xFF, 0xD1, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xB0, 0xB8, 0xD0)
MGRAY     = RGBColor(0x60, 0x6B, 0x8A)
DGRAY     = RGBColor(0x10, 0x14, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 基础工具 ───────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def box(s, l, t, w, h, color):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def tb(s, text, l, t, w, h, size=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    txb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def hline(s, t, color=BLUE, l=0.5, w=12.33):
    box(s, l, t, w, 0.03, color)

def page_footer(s):
    box(s, 0, 7.18, 13.33, 0.32, DGRAY)
    tb(s, "GOBS v2.0  ·  2026年4月  ·  内部机密", 0.5, 7.21, 12.33, 0.28,
       size=9, color=MGRAY, align=PP_ALIGN.CENTER)

def top_bar(s, color=BLUE):
    box(s, 0, 0, 13.33, 0.07, color)

def section_tag(s, text, color=BLUE):
    box(s, 0.5, 0.18, 0.07, 0.3, color)
    tb(s, text, 0.65, 0.18, 8, 0.3, size=9, bold=True, color=color)

def slide_heading(s, title, sub=None):
    tb(s, title, 0.5, 0.52, 12.33, 0.55, size=26, bold=True, color=WHITE)
    if sub:
        tb(s, sub, 0.5, 1.1, 12.33, 0.32, size=12, color=LGRAY)
    hline(s, 1.48)

def card(s, l, t, w, h, color=CARD):
    return box(s, l, t, w, h, color)


# ══════════════════════════════════════════════════════════════════════════════
# 第1页 — 封面
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 0.07, BLUE)

# 左侧主体
tb(s, "GOBS", 0.6, 1.6, 8, 1.5, size=90, bold=True, color=WHITE)
box(s, 0.6, 3.1, 4.5, 0.07, BLUE)
tb(s, "游戏运营与内容发行系统", 0.62, 3.25, 9, 0.5, size=20, bold=True, color=WHITE)
tb(s, "Game Operations & Broadcasting System",
   0.62, 3.82, 9, 0.35, size=13, color=LGRAY)
tb(s, "产品愿景  ·  架构规划  ·  落地路线图",
   0.62, 4.22, 9, 0.35, size=12, color=MGRAY)
box(s, 0.62, 4.65, 1.8, 0.32, PURPLE)
tb(s, "v2.0  |  2026年4月", 0.72, 4.68, 1.6, 0.28, size=10, bold=True, color=WHITE)

# 右侧装饰
box(s, 9.2, 1.2, 3.8, 5.2, CARD)
box(s, 9.2, 1.2, 3.8, 0.06, PURPLE)
tb(s, "🚀", 10.3, 2.6, 1.6, 1.2, size=60, align=PP_ALIGN.CENTER)
tb(s, "智能发行飞轮", 9.4, 4.1, 3.4, 0.4, size=14, bold=True,
   color=PURPLE, align=PP_ALIGN.CENTER)
tb(s, "感知 · 洞察 · 执行 · 学习",
   9.4, 4.58, 3.4, 0.35, size=11, color=LGRAY, align=PP_ALIGN.CENTER)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第2页 — 目录
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s)
section_tag(s, "目录")
slide_heading(s, "今日议程")

items = [
    ("01", "我们在哪里",    "梳理 GOBS H5 当前已有能力与现存缺口",         BLUE),
    ("02", "我们要去哪里",  "愿景：AI 驱动的游戏发行飞轮",                PURPLE),
    ("03", "如何实现",      "架构分层、核心模块拆解与时间节点规划",        GREEN),
]
for i, (num, title, desc, color) in enumerate(items):
    t0 = 1.65 + i * 1.7
    card(s, 0.5, t0, 12.33, 1.55)
    box(s, 0.5, t0, 0.08, 1.55, color)
    tb(s, num,   0.75, t0 + 0.2,  0.9,  0.5,  size=28, bold=True, color=color)
    tb(s, title, 1.78, t0 + 0.18, 5,    0.45, size=20, bold=True, color=WHITE)
    tb(s, desc,  1.78, t0 + 0.72, 10.5, 0.5,  size=12, color=LGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第3页 — Section 01 分隔页
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 0.12, 7.5, BLUE)
tb(s, "01", 1.0, 1.2, 5, 1.8, size=110, bold=True, color=RGBColor(0x1A,0x20,0x35))
tb(s, "我们在哪里", 1.0, 2.5, 10, 0.95, size=52, bold=True, color=WHITE)
box(s, 1.0, 3.55, 4.5, 0.06, BLUE)
tb(s, "WHERE WE ARE", 1.0, 3.68, 6, 0.38, size=14, color=BLUE, bold=True)
tb(s, "盘点 GOBS H5 当前的核心能力，以及飞轮尚未闭合的关键缺口。",
   1.0, 4.18, 10, 0.45, size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 第4页 — 当前平台能力
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s)
section_tag(s, "01  我们在哪里")
slide_heading(s, "GOBS H5 — 当前平台能力", "今天已经可以做的事")

modules = [
    ("🎬", "内容生产", [
        "Studio：文本/参考图驱动生视频",
        "Seedance 2.0 集成（ViralDance）",
        "多分镜故事板 + 全局风格锁定",
        "角色形象编辑器 + 形象演化树",
    ], BLUE),
    ("🧑‍🎨", "角色系统", [
        "角色卡 + 形象树（多变体管理）",
        "状态衣橱：基础形象 + 多状态变体",
        "形象库：保存并跨项目复用",
        "场景道具卡生图",
    ], PURPLE),
    ("📤", "分发系统", [
        "GeeLark 账号矩阵集成",
        "夜间批量定时发布",
        "TikTok / Instagram 多账号支持",
        "任务看板 + 状态追踪",
    ], GREEN),
    ("🗂️", "高级制片", [
        "完整制片向导流程",
        "模板市场 + 保存为模板",
        "Drive 素材选取器",
        "批量任务队列管理",
    ], ORANGE),
]

# 2×2 布局，行间距足够
for i, (icon, title, bullets, color) in enumerate(modules):
    col = i % 2; row = i // 2
    l = 0.5 + col * 6.45
    t = 1.65 + row * 2.8
    card(s, l, t, 6.2, 2.65)
    box(s, l, t, 6.2, 0.06, color)
    tb(s, icon + "  " + title, l + 0.2, t + 0.15, 5.8, 0.42,
       size=13, bold=True, color=color)
    for j, b in enumerate(bullets):
        tb(s, "·  " + b, l + 0.25, t + 0.66 + j * 0.47, 5.7, 0.42, size=10.5, color=LGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第5页 — 当前缺口
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s)
section_tag(s, "01  我们在哪里  ·  缺口分析")
slide_heading(s, "飞轮断在哪里？", "生产与分发之后，什么都没发生")

# 流程图
steps = [
    ("内容\n生产", BLUE,   True),
    ("分发\n发布", PURPLE, True),
    ("表现\n数据", ORANGE, False),
    ("洞察\n决策", GREEN,  False),
]
sx = [0.55, 3.55, 6.55, 9.55]
for i, ((label, color, exists), x) in enumerate(zip(steps, sx)):
    bg_c = CARD if exists else RGBColor(0x22, 0x10, 0x10)
    card(s, x, 1.65, 2.7, 1.45, color=bg_c)
    box(s, x, 1.65, 2.7, 0.06, color if exists else RED)
    if not exists:
        tb(s, "❌ 缺失", x + 0.15, 1.72, 1.5, 0.28, size=9, bold=True, color=RED)
    tb(s, label, x + 0.15, 1.95 if not exists else 1.85,
       2.4, 0.75, size=14, bold=exists, color=WHITE if exists else MGRAY,
       align=PP_ALIGN.CENTER)
    if i < 3:
        tb(s, "→", x + 2.72, 2.15, 0.35, 0.5, size=18, bold=True,
           color=MGRAY, align=PP_ALIGN.CENTER)

# 4个缺口说明
gaps = [
    ("无数据回流",    "发布后视频表现数据（播放量/完播率/互动）从未回到平台。「做了什么」和「结果如何」之间没有因果链路。", RED),
    ("无智能洞察",    "零 AI 归因分析。不知道哪类内容、哪个角色、哪个时间段表现最好。所有决策靠经验和感觉。",            ORANGE),
    ("无行动建议",    "即使有数据，也没有机制把信号转化成具体 Action：做这条视频、联系这个 KOL、修复这个游戏系统。",      YELLOW),
    ("无自我学习",    "平台无法从历史 Action 中学习。每次 campaign 从零开始，没有复利积累，没有智能进化。",               PURPLE),
]
for i, (title, desc, color) in enumerate(gaps):
    col = i % 2; row = i // 2
    l = 0.5 + col * 6.45
    t = 3.3 + row * 1.95
    card(s, l, t, 6.2, 1.85)
    box(s, l, t, 0.07, 1.85, color)
    tb(s, title, l + 0.22, t + 0.15, 5.8, 0.38, size=13, bold=True, color=color)
    tb(s, desc,  l + 0.22, t + 0.58, 5.8, 1.1,  size=10.5, color=LGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第6页 — Section 02 分隔页
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 0.12, 7.5, PURPLE)
tb(s, "02", 1.0, 1.2, 5, 1.8, size=110, bold=True, color=RGBColor(0x1A,0x20,0x35))
tb(s, "我们要去哪里", 1.0, 2.5, 10, 0.95, size=46, bold=True, color=WHITE)
box(s, 1.0, 3.55, 5, 0.06, PURPLE)
tb(s, "WHERE WE AIM TO GO", 1.0, 3.68, 7, 0.38, size=14, color=PURPLE, bold=True)
tb(s, "AI 驱动的智能发行飞轮——感知、洞察、执行、学习，最终自主运转。",
   1.0, 4.18, 10, 0.45, size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 第7页 — 愿景与飞轮
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, PURPLE)
section_tag(s, "02  我们要去哪里", color=PURPLE)
slide_heading(s, "GOBS 愿景", "绑定游戏，AI 接管发行运转")

# 核心理念
tb(s, "用户绑定游戏后，平台持续采集内外部信号，AI 生成基于发行阶段和 ROI 的行动建议，"
      "Agent 自动执行并追踪效果，系统不断自我学习，最终以小团队驾驭大发行规模。",
   0.5, 1.6, 12.33, 0.65, size=12, color=LGRAY)

# 飞轮 4 个节点
flywheel = [
    ("📡", "感知",  "游戏KPI · 内容表现\n买量数据 · 舆情 · KOL",  BLUE,   0.5),
    ("🧠", "洞察",  "AI归因 · 阶段感知\nROI排序行动建议",          PURPLE, 3.65),
    ("⚡", "执行",  "视频生产 · 分发\nKOL · 危机 · PRD",           ORANGE, 6.8),
    ("📊", "追踪",  "ROI归因\n行动→效果因果记录",                  GREEN,  9.95),
]
for icon, title, desc, color, l in flywheel:
    card(s, l, 2.4, 3.0, 3.1)
    box(s, l, 2.4, 3.0, 0.06, color)
    tb(s, icon, l + 1.1, 2.55, 0.85, 0.7, size=30, align=PP_ALIGN.CENTER)
    tb(s, title, l + 0.1, 3.32, 2.8, 0.45, size=16, bold=True,
       color=color, align=PP_ALIGN.CENTER)
    tb(s, desc, l + 0.1, 3.85, 2.8, 0.9, size=10.5, color=LGRAY, align=PP_ALIGN.CENTER)

# 箭头
for ax in [3.62, 6.77, 9.92]:
    tb(s, "→", ax, 3.65, 0.2, 0.5, size=18, bold=True,
       color=PURPLE, align=PP_ALIGN.CENTER)

tb(s, "↩  自学习反馈闭环", 4.3, 5.62, 4.5, 0.38,
   size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# 自动化三阶段
modes = [
    ("辅助模式",   "AI 推荐，人工确认每个 Action",     "0–3 个月",   BLUE),
    ("半自动模式", "低风险 Action 自动执行，高风险人审", "3–6 个月",  PURPLE),
    ("全自动模式", "设定方向与红线，平台自主运转",       "6 个月以上", GREEN),
]
for i, (m, d, when, color) in enumerate(modes):
    l = 0.5 + i * 4.27
    card(s, l, 6.05, 4.0, 1.12)
    box(s, l, 6.05, 4.0, 0.06, color)
    tb(s, m,    l + 0.18, 6.15, 3.6, 0.35, size=12, bold=True, color=color)
    tb(s, d,    l + 0.18, 6.55, 3.6, 0.38, size=10, color=LGRAY)
    tb(s, when, l + 0.18, 6.96, 3.6, 0.2,  size=9,  color=MGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第8页 — 完整架构分层
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, PURPLE)
section_tag(s, "02  我们要去哪里  ·  架构", color=PURPLE)
slide_heading(s, "完整平台架构（8层）", "从数据采集到自主执行的全栈设计")

layers = [
    ("L0", "接入层",         "游戏绑定 · 发行阶段配置 · 合规设置 · 知识库初始化",                             BLUE),
    ("L1", "数据感知层",     "游戏内KPI · 内容表现 · 买量数据(Meta/TikTok Ads) · 舆情 · KOL · 竞品动向",     GREEN),
    ("L2", "洞察引擎",       "归因分析 · 阶段感知建议 · ROI权重排序 · 每日行动清单",                          PURPLE),
    ("L3", "合规检查层",     "内容规则 · 地区过滤 · 风险评级（自动放行/人工审核/拦截）",                      ORANGE),
    ("L4", "Action Agent层", "内容Agent · 分发Agent · KOL Agent · 危机Agent · 产品Agent · 买量Agent",        BLUE),
    ("L5", "人工审核门",     "官号发布 · KOL合作确认 · 危机公关 · 预算调整  →  永远需要人工",                RED),
    ("L6", "ROI追踪层",      "每次Action成本 · 增量价值归因 · 周度ROI排行 · 预算分配建议",                    YELLOW),
    ("L7", "记忆与自学习层", "行动档案（触发→执行→结果） · 策略权重更新 · 跨游戏经验迁移",                   PURPLE),
]

row_h = 0.65
for i, (code, name, desc, color) in enumerate(layers):
    t = 1.58 + i * row_h
    bg_c = CARD if i % 2 == 0 else CARD2
    card(s, 0.5, t, 12.33, row_h - 0.04, color=bg_c)
    box(s, 0.5, t, 0.08, row_h - 0.04, color)
    tb(s, code, 0.65, t + 0.09, 0.55, 0.32, size=8.5, bold=True, color=color)
    tb(s, name, 1.26, t + 0.05, 2.2,  0.32, size=12,  bold=True, color=WHITE)
    tb(s, desc, 3.55, t + 0.09, 9.1,  0.45, size=10,  color=LGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第9页 — Section 03 分隔页
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 0.12, 7.5, GREEN)
tb(s, "03", 1.0, 1.2, 5, 1.8, size=110, bold=True, color=RGBColor(0x1A,0x20,0x35))
tb(s, "如何实现", 1.0, 2.5, 10, 0.9, size=52, bold=True, color=WHITE)
box(s, 1.0, 3.5, 4, 0.06, GREEN)
tb(s, "HOW TO ACHIEVE", 1.0, 3.63, 7, 0.38, size=14, color=GREEN, bold=True)
tb(s, "核心模块拆解、功能点清单，以及有明确时间节点的交付路线图。",
   1.0, 4.13, 10, 0.45, size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 第10页 — Phase 1：地基（第1个月）
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, GREEN)
section_tag(s, "03  如何实现  ·  Phase 1", color=GREEN)
slide_heading(s, "Phase 1 — 地基", "第1个月  |  目标：打通数据闭环")

# 标签
box(s, 0.5, 1.5, 1.5, 0.3, GREEN)
tb(s, "第1个月", 0.55, 1.52, 1.42, 0.26, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cols = [
    ("游戏接入模块", [
        ("游戏绑定：名称/地区/平台/团队配置",         "第1周"),
        ("发行阶段配置（软上线/爆发期/长线运营）",    "第1周"),
        ("官号 OAuth 授权（TikTok/YouTube/IG）",      "第2周"),
        ("合规配置：地区分级 + 内容红线设置",         "第2周"),
        ("知识库初始化（Wiki/角色设定/历史内容）",    "第3周"),
    ], BLUE),
    ("数据回流（L1 MVP）", [
        ("官号内容表现数据接入（平台API）",            "第2周"),
        ("视频→效果因果记录（模板/角色/时间/平台）",  "第2周"),
        ("每日表现数据 Dashboard",                    "第3周"),
        ("游戏内数据接入（DAU/MAU，API或手动）",      "第3–4周"),
        ("买量数据接入（TikTok Ads / Meta Ads API）", "第4周"),
    ], GREEN),
    ("内容生产优化", [
        ("角色卡「保存到形象库」入口  ✅ 已完成",     "Done"),
        ("状态衣橱基础形象自动带入定稿图  ✅ 已完成", "Done"),
        ("批量变体生产（同方向3条，A/B测试）",        "第3周"),
        ("品质自动评分（完整度/动作流畅度）",         "第4周"),
        ("高表现视频自动沉淀为模板",                  "第4周"),
    ], PURPLE),
]

col_l = [0.5, 4.62, 8.74]
for ci, (title, items, color) in enumerate(cols):
    l = col_l[ci]
    card(s, l, 1.9, 4.0, 5.2)
    box(s, l, 1.9, 4.0, 0.06, color)
    tb(s, title, l + 0.15, 2.02, 3.7, 0.38, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        done = "Done" in when
        t0 = 2.52 + j * 0.93
        fc = GREEN if done else LGRAY
        wc = GREEN if done else MGRAY
        tb(s, "✓" if done else "·", l + 0.12, t0,        0.22, 0.35, size=10, color=fc)
        tb(s, feat,                  l + 0.33, t0,        3.5,  0.45, size=10, color=fc)
        tb(s, when,                  l + 0.33, t0 + 0.45, 3.5,  0.25, size=8,  color=wc)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第11页 — Phase 2：闭环（第2–3个月）
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, PURPLE)
section_tag(s, "03  如何实现  ·  Phase 2", color=PURPLE)
slide_heading(s, "Phase 2 — 闭环", "第2–3个月  |  目标：洞察驱动行动")

box(s, 0.5, 1.5, 1.8, 0.3, PURPLE)
tb(s, "第2–3个月", 0.55, 1.52, 1.72, 0.26, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cols = [
    ("洞察引擎（L2）", [
        ("舆情数据接入（Reddit/TikTok/Twitter情绪分析）", "第2个月"),
        ("内容归因：角色/模板/时间段效果分析",            "第2个月"),
        ("阶段感知：同信号在不同发行阶段给不同建议",      "第2个月"),
        ("竞品监控：内容趋势/买量动向/版本节奏",          "第2–3个月"),
        ("ROI基准：每类Action成本 vs 增量价值",           "第3个月"),
        ("每日洞察报告 + 优先级排序行动清单",             "第3个月"),
    ], PURPLE),
    ("行动推荐层", [
        ("行动Brief生成（触发原因/执行内容/预期效果）",  "第2个月"),
        ("人工确认流程（确认→队列→Agent执行）",         "第2个月"),
        ("KOL数据库：分级档案（一/二/三线）/受众/报价", "第2–3个月"),
        ("KOL合作机会识别 + 邀约文案起草",              "第3个月"),
        ("危机检测：负面情绪触发阈值→进入危机模式",     "第3个月"),
        ("危机模式：暂停所有发布，辅助起草官方回应",    "第3个月"),
    ], ORANGE),
    ("分发与合规", [
        ("地区内容变体（字幕/配音按地区自动生成）",  "第2个月"),
        ("合规门：每次发布前自动过规则检查",         "第2个月"),
        ("风险评级：自动放行/标记/拦截三档",         "第2个月"),
        ("最优发布时间推荐（基于历史数据ML）",       "第3个月"),
        ("多平台同步发布 + 平台格式适配",            "第3个月"),
    ], GREEN),
]

col_l = [0.5, 4.62, 8.74]
for ci, (title, items, color) in enumerate(cols):
    l = col_l[ci]
    card(s, l, 1.9, 4.0, 5.2)
    box(s, l, 1.9, 4.0, 0.06, color)
    tb(s, title, l + 0.15, 2.02, 3.7, 0.38, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        t0 = 2.52 + j * 0.78
        tb(s, "·", l + 0.12, t0, 0.22, 0.35, size=10, color=LGRAY)
        tb(s, feat, l + 0.3,  t0, 3.55, 0.42, size=9.5, color=LGRAY)
        tb(s, when, l + 0.3,  t0 + 0.43, 3.55, 0.24, size=8, color=MGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第12页 — Phase 3：智能化（第4–6个月）
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, ORANGE)
section_tag(s, "03  如何实现  ·  Phase 3", color=ORANGE)
slide_heading(s, "Phase 3 — 智能化", "第4–6个月  |  目标：自学习，半自动运转")

box(s, 0.5, 1.5, 1.8, 0.3, ORANGE)
tb(s, "第4–6个月", 0.55, 1.52, 1.72, 0.26, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cols = [
    ("全套 Action Agent（L4）", [
        ("内容Agent：prompt生成→Seedance→品质评分→审核队列", "第4个月"),
        ("分发Agent：最优时间自动发布，官号+矩阵号",          "第4个月"),
        ("KOL Agent：合作追踪/邀约起草/效果归因",             "第4个月"),
        ("买量Agent：高效素材特征提取/预算再分配建议",        "第4–5个月"),
        ("产品Agent：玩家反馈→根因分析→PRD草稿→Jira",        "第5个月"),
        ("危机Agent：舆情监控→危机模式→回应起草",            "第5个月"),
    ], ORANGE),
    ("记忆与自学习（L7）", [
        ("行动档案：触发→执行→ROI结果完整因果记录",    "第4个月"),
        ("周度策略权重更新：高ROI行动类型优先推荐",    "第5个月"),
        ("Agent默认参数调优（时间/风格/角色偏好）",    "第5个月"),
        ("游戏新版本发布→知识库自动更新",              "第5个月"),
        ("负反馈学习：被拒绝Action降低推荐权重",       "第5–6个月"),
        ("行业Benchmark数据库（各品类正常完播率/CPM）","第6个月"),
    ], PURPLE),
    ("ROI追踪与汇报（L6）", [
        ("每次Action成本模型（工具+人力当量估算）",     "第4个月"),
        ("增量价值归因（DAU/营收提升归因到具体Action）","第4–5个月"),
        ("周度ROI排行榜：哪类Action性价比最高",         "第5个月"),
        ("预算分配建议：把钱花在ROI最高的Action上",    "第5个月"),
        ("半自动模式上线：低风险Action免审核自动执行", "第6个月"),
        ("跨游戏经验迁移：A游戏有效策略推荐给B游戏",  "第6个月"),
    ], GREEN),
]

col_l = [0.5, 4.62, 8.74]
for ci, (title, items, color) in enumerate(cols):
    l = col_l[ci]
    card(s, l, 1.9, 4.0, 5.2)
    box(s, l, 1.9, 4.0, 0.06, color)
    tb(s, title, l + 0.15, 2.02, 3.7, 0.38, size=12, bold=True, color=color)
    for j, (feat, when) in enumerate(items):
        t0 = 2.52 + j * 0.78
        tb(s, "·", l + 0.12, t0, 0.22, 0.35, size=10, color=LGRAY)
        tb(s, feat, l + 0.3,  t0, 3.55, 0.42, size=9.5, color=LGRAY)
        tb(s, when, l + 0.3,  t0 + 0.43, 3.55, 0.24, size=8, color=MGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第13页 — Phase 4：飞轮（6个月以上）
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, GREEN)
section_tag(s, "03  如何实现  ·  Phase 4", color=GREEN)
slide_heading(s, "Phase 4 — 飞轮", "6个月以上  |  目标：平台自主运转")

box(s, 0.5, 1.5, 2.0, 0.3, GREEN)
tb(s, "6个月以上", 0.55, 1.52, 1.9, 0.26, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

items = [
    ("全自动模式",
     "低风险 Action 自动执行，无需每次审批。人工设定方向与红线，周度复盘替代日常审核。",
     GREEN, 0.5, 1.95),
    ("三维优化矩阵",
     "账号画像 × 内容类型 × 发布时间，由学习层持续优化，覆盖各平台各地区。",
     PURPLE, 0.5, 3.45),
    ("跨游戏智能迁移",
     "游戏A的有效策略为游戏B提供参考。平台逐步积累各品类的发行 Playbook。",
     ORANGE, 0.5, 4.95),
    ("智能发行飞轮闭环",
     "感知→洞察→执行→追踪→学习→循环。人的角色从「操盘手」变为「方向定义者」。10人团队发挥100人规模。",
     BLUE, 6.95, 1.95),
    ("平台级行业智能",
     "跨游戏聚合数据（脱敏）形成行业基准。GOBS 成为该品类「好的发行是什么样子」的参照标准。",
     PURPLE, 6.95, 3.45),
    ("永远保留人工的边界",
     "品牌官号发布 · KOL合同签署 · 危机公关最终定稿 · 重大预算决策 · 产品路线方向",
     RED, 6.95, 4.95),
]
for title, desc, color, l, t in items:
    card(s, l, t, 5.95, 1.3)
    box(s, l, t, 0.08, 1.3, color)
    tb(s, title, l + 0.22, t + 0.1,  5.55, 0.38, size=13, bold=True, color=color)
    tb(s, desc,  l + 0.22, t + 0.52, 5.55, 0.72, size=10.5, color=LGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第14页 — 总路线图
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s)
section_tag(s, "03  如何实现  ·  路线图总览")
slide_heading(s, "交付路线图", "各阶段核心模块交付时间线")

# 阶段表头
phases = [
    ("Phase 1\n第1个月",    BLUE,   0.5),
    ("Phase 2\n第2–3个月",  PURPLE, 3.82),
    ("Phase 3\n第4–6个月",  ORANGE, 7.14),
    ("Phase 4\n6个月以上",  GREEN,  10.46),
]
for label, color, l in phases:
    box(s, l, 1.55, 3.08, 0.58, color)
    tb(s, label, l + 0.08, 1.57, 2.9, 0.55, size=9.5, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

# 行数据
rows = [
    ("数据层",       ["游戏绑定+知识库",       "舆情+竞品接入",          "买量深度分析",           "跨游戏Benchmark"]),
    ("洞察引擎",     ["表现数据Dashboard",     "归因+阶段感知建议",      "ROI排序行动",            "三维优化矩阵"]),
    ("合规检查",     ["合规配置搭建",          "合规门+风险评级",        "—",                      "—"]),
    ("Action Agent", ["内容Agent MVP",          "分发+KOL Agent",         "全套Agent+危机",         "全自动模式"]),
    ("人工审核门",   ["所有Action",            "仅高风险",               "仅高风险",               "品牌/预算/危机"]),
    ("ROI追踪",      ["—",                     "成本基准模型",           "ROI排行+预算建议",       "—"]),
    ("学习层",       ["—",                     "—",                      "行动档案+权重更新",      "跨游戏迁移"]),
]
col_colors = [BLUE, PURPLE, ORANGE, GREEN]
row_h = 0.67
for ri, (module, cells) in enumerate(rows):
    t = 2.2 + ri * row_h
    bg_c = CARD if ri % 2 == 0 else CARD2
    # 模块名
    box(s, 0.0, t, 0.5, row_h, DGRAY)
    tb(s, module, 0.02, t + 0.18, 0.48, 0.35, size=7, bold=True, color=BLUE)
    for ci, cell in enumerate(cells):
        l = 0.5 + ci * 3.32
        box(s, l, t, 3.28, row_h, bg_c)
        fc = WHITE if ci >= 2 else LGRAY
        if cell != "—":
            tb(s, "· " + cell, l + 0.1, t + 0.15, 3.05, 0.48, size=9.5, color=fc)
        else:
            tb(s, "—", l + 1.3, t + 0.2, 0.5, 0.3, size=9, color=MGRAY, align=PP_ALIGN.CENTER)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第15页 — 自动化边界
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s, RED)
section_tag(s, "03  如何实现  ·  自动化边界", color=RED)
slide_heading(s, "自动化 vs 人工判断", "知道哪些不能自动化，和知道哪些可以，同样重要")

# 左：可全自动
card(s, 0.5, 1.6, 5.9, 5.7)
box(s, 0.5, 1.6, 5.9, 0.06, GREEN)
tb(s, "✅  可以全自动", 0.65, 1.72, 5.5, 0.42, size=14, bold=True, color=GREEN)

auto_items = [
    "数据采集与每日 Dashboard 更新",
    "归因分析与洞察报告生成",
    "视频生产（进审核队列，不直接发布）",
    "矩阵号/小号内容分发",
    "内部洞察报告与 PRD 草稿生成",
    "品质评分不达标视频自动重生",
    "竞品与舆情持续监控",
    "最优发布时间推荐",
]
for i, item in enumerate(auto_items):
    tb(s, "·  " + item, 0.65, 2.22 + i * 0.62, 5.6, 0.55, size=11, color=LGRAY)

# 右：必须人工
card(s, 6.93, 1.6, 5.9, 5.7)
box(s, 6.93, 1.6, 5.9, 0.06, RED)
tb(s, "🚫  必须人工审核", 7.08, 1.72, 5.5, 0.42, size=14, bold=True, color=RED)

human_items = [
    ("品牌官号公开内容发布",     "错误成本高，一旦发出影响品牌形象"),
    ("KOL合作确认与签约",        "商务关系，需要人工判断和谈判"),
    ("危机公关最终回应",         "声誉攸关，措辞错误会扩大危机"),
    ("跨地区敏感内容发布",       "文化差异，AI无法完全判断"),
    ("重大预算调整（超过阈值）", "财务风险，需要明确责任主体"),
    ("产品PRD正式提交",          "业务影响，需要利益相关方对齐"),
]
for i, (action, reason) in enumerate(human_items):
    t0 = 2.22 + i * 0.85
    tb(s, "⚠  " + action, 7.08, t0,        5.6, 0.38, size=11, bold=True, color=WHITE)
    tb(s, reason,          7.08, t0 + 0.4,  5.6, 0.38, size=9.5, color=MGRAY)

page_footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# 第16页 — 结语
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); top_bar(s)

tb(s, "飞轮的启动\n只需要一个决定：", 0.8, 1.0, 10, 1.4,
   size=30, bold=True, color=WHITE)

tb(s, "打通数据闭环。", 0.8, 2.5, 10, 0.85, size=46, bold=True, color=BLUE)

box(s, 0.8, 3.45, 5.5, 0.06, BLUE)

tb(s, "感知→洞察→执行→追踪→学习——这套飞轮中的一切，"
      "都建立在「我们做了什么」和「结果怎么样」能被连接的前提上。"
      "Phase 1 是地基，先把它做扎实。",
   0.8, 3.6, 11.5, 0.75, size=12.5, color=LGRAY)

priorities = [
    ("第1步", "数据回流\n（内容→表现）",    BLUE),
    ("第2步", "洞察引擎MVP\n（阶段感知）",   PURPLE),
    ("第3步", "Action Agent\n+合规门",        GREEN),
    ("第4步", "自学习\n+ROI追踪",            ORANGE),
]
for i, (rank, desc, color) in enumerate(priorities):
    l = 0.8 + i * 3.0
    card(s, l, 4.6, 2.75, 1.72)
    box(s, l, 4.6, 2.75, 0.06, color)
    tb(s, rank, l + 0.15, 4.72, 2.45, 0.42, size=13, bold=True, color=color)
    tb(s, desc, l + 0.15, 5.18, 2.45, 0.95, size=11, color=LGRAY)

tb(s, "10人团队  ×  GOBS  =  100人规模的发行能力",
   0.8, 6.55, 11.5, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

page_footer(s)


# ── 保存 ──────────────────────────────────────────────────────────────────────
out = "/home/openclaw/.openclaw/workspace/gobs_review/GOBS_平台规划_v2.pptx"
prs.save(out)
print(f"已保存：{out}")
print(f"共 {len(prs.slides)} 页")
